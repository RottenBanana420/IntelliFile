"""
Import necessary libraries for text extraction and processing.
"""

import json
import csv
import os

# Import PDFMiner library for extracting text from PDF files
import pdfminer
from pdfminer.high_level import extract_text
from pdfminer.layout import LAParams

# Import eBookLib library for extracting text from EPUB files
import ebooklib
from ebooklib import epub

# Import python-docx library for extracting text from DOCX files
import docx

# Import pandas library for data manipulation and analysis
import pandas as pd

# Import python-pptx library for extracting text from PPTX files
import pptx

# Load environment variables from .env file
from dotenv import load_dotenv
load_dotenv()

# Import LangChain library for language model interactions
from langchain_community.llms import Ollama
from langchain_core.output_parsers import StrOutputParser

# Suppress warnings to avoid cluttering the output
import warnings
warnings.filterwarnings("ignore")
# Get the categories from the environment variable 'CATEGORIES'
categories_str = os.getenv('CATEGORIES')

# Split the categories string into a list using comma as the separator
categories_list = categories_str.split(',')
def read_file(file_path, mode='r'):
    """
    Reads the content of a file based on its extension.

    Supported file types:
    - .txt: Plain text file
    - .json: JSON file
    - .xlsx: Excel file
    - .csv: CSV file
    - .py, .java, .cpp, .c, .js, .ts, .go, .rb, .swift, .kt, .scala, .php, .perl, .ruby, .bash, .sh, .zsh, .html: Code files
    - .epub: EPUB file
    - .pdf: PDF file
    - .docx: Word document
    - .ppt, .pptx: PowerPoint presentation

    Args:
        file_path (str): The path to the file.
        mode (str, optional): The mode in which the file is opened. Default is 'r'.

    Returns:
        content: The content of the file. Type depends on the file extension.

    Raises:
        FileNotFoundError: If the file does not exist at the given path.
    """
    try:
        if file_path.endswith('.txt'):
            # Read plain text file
            with open(file_path, mode) as file:
                content = file.read()
                return content

        elif file_path.endswith('.json'):
            # Read JSON file
            with open(file_path, mode) as file:
                content = json.load(file)
                return content

        elif file_path.endswith('.xlsx'):
            # Read Excel file
            excel_file = pd.read_excel(file_path)
            content = excel_file.to_string()
            return content

        elif file_path.endswith('.csv'):
            # Read CSV file
            with open(file_path, mode) as file:
                reader = csv.reader(file)
                content = list(reader)
                return content
            
        elif file_path.endswith(('.py', '.java', '.cpp', '.c', '.js', '.ts', '.go', '.rb', '.swift', '.kt', '.scala', '.php', '.perl', '.ruby', '.bash', '.sh', '.zsh', '.html')):
            # Read code file
            with open(file_path, mode) as file:
                content = file.read()
                return content

        elif file_path.endswith('.epub'):
            # Read EPUB file
            book = epub.read_epub(file_path)
            content = ''
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                content += item.get_content().decode('utf-8')
            return content

        elif file_path.endswith('.pdf'):
            # Read PDF file
            content = extract_text(file_path)
            return content
            
        elif file_path.endswith('.docx'):
            # Read Word document
            doc = docx.Document(file_path)
            content = ''
            for para in doc.paragraphs:
                content += para.text
            return content
   
        elif file_path.endswith(('.ppt', '.pptx')):
            # Read PowerPoint presentation
            presentation = pptx.Presentation(file_path)
            content = ''
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        content += shape.text_frame.text
            return content

    except FileNotFoundError:
        print(f"File not found: {file_path}")
        return None
def rename_file(file_path: str, new_name: str, category: str = None) -> None:
    """
    Renames a file by appending a category and a new name to its existing extension.
    
    If category is not provided, the file will be renamed with just the new name.

    Args:
        file_path (str): The path to the file to be renamed.
        new_name (str): The new name to be given to the file.
        category (str, optional): The category to prepend to the new file name. Defaults to None.

    Returns:
        None

    Raises:
        FileNotFoundError: If the file does not exist at the given path.
        PermissionError: If there is no permission to rename the file.
    """
    # Split the file path into directory and file name
    file_dir, file_name = os.path.split(file_path)
    # Get the file extension
    file_ext = os.path.splitext(file_name)[1]
    
    if category is not None:
        # Create the new file path with category and new name
        new_file_path = os.path.join(file_dir, f"{category}_{new_name}{file_ext}")
    else:
        # Create the new file path with just the new name
        new_file_path = os.path.join(file_dir, f"{new_name}{file_ext}")

    try:
        # Attempt to rename the file
        os.rename(file_path, new_file_path)
    except FileNotFoundError:
        print(f"Error: The file '{file_path}' does not exist.")
    except PermissionError:
        print(f"Error: Permission denied to rename the file '{file_path}'.")
# def main():
#     """
#     Renames files in a specified folder based on their content.

#     The function reads the content of each file, generates a descriptive and unique name for it,
#     assigns a category to the file, and renames the file accordingly.

#     Folder and file processing:
#     - Skips hidden files.
#     - Reads the content of each file.
#     - Generates a new file name using an LLM model.
#     - Assigns a category using the LLM model.
#     - Renames the file with the new name and category.

#     Args:
#         None

#     Returns:
#         None
#     """
#     folder_path = 'data'  # Define the folder path containing the files
    
#     for file_name in os.listdir(folder_path):
#         if not file_name.startswith('.'):
#             file_path = os.path.join(folder_path, file_name)
#             if os.path.isfile(file_path):
#                 # Read the file content
#                 content = read_file(file_path)
                
#                 llm = Ollama(model="llama3.1")

#                 # Generate a descriptive and unique name for the file
#                 prompt1 = (f"Generate a descriptive and unique name for a file based on the following content: {content}. "
#                            f"The name should be concise, informative, and relevant to the content. Keep the name small, limit it to 3 words. "
#                            f"Please provide just the name, with a single option that you think is best, without any file extensions. Respond with just the name and nothing else.")
#                 response = llm(prompt1)
#                 response = response.strip('\"')
                
#                 # Assign a general category to the file
#                 prompt2 = (f"Assign a general category from the predefined list {categories_list} to a file based on its content {content}. "
#                            f"The category should be one of the following: {categories_list}. Do not assign a category that is not in this list. "
#                            f"Respond with just the category and nothing else.")
#                 category = llm(prompt2)
                
#                 # Rename the file
#                 rename_file(file_path, response, category)
#                 print(f"Renamed {file_name} to {category}_{response}")
def main():
    """
    Renames files in a specified folder based on their content.

    The function reads the content of each file, generates a descriptive and unique name for it,
    assigns a category to the file, and renames the file accordingly.

    Folder and file processing:
    - Skips hidden files.
    - Reads the content of each file.
    - Generates a new file name using an LLM model.
    - Assigns a category using the LLM model.
    - Renames the file with the new name and category.

    Args:
        None

    Returns:
        None
    """
    folder_path = 'data'  # Define the folder path containing the files

    for file_name in os.listdir(folder_path):
        if not file_name.startswith('.'):
            file_path = os.path.join(folder_path, file_name)
            if os.path.isfile(file_path):
                # Read the file content
                content = read_file(file_path)
                
                llm = Ollama(model="llama3.1")

                # Generate a descriptive and unique name for the file
                prompt1 = (f"Generate a descriptive and unique name for a file based on the following content: {content}. "
                           f"The name should be concise, informative, and relevant to the content. Keep the name small, limit it to 3 words. "
                           f"Please provide just the name, with a single option that you think is best, without any file extensions. Respond with just the name and nothing else.")
                response = llm(prompt1)
                response = response.strip('\"')
                
                # Assign a general category to the file
                prompt2 = (f"Assign a general category from the predefined list {categories_list} to a file based on its content {content}. "
                           f"The category should be one of the following: {categories_list}. Do not assign a category that is not in this list. "
                           f"Respond with just the category and nothing else.")
                category = llm(prompt2)
                
                # Check if the generated name is acceptable (less than 5 words)
                if len(response.split()) < 5:
                    # Check if the generated category is in the predefined list
                    if category in categories_list:
                        # Rename the file with the new name and category
                        rename_file(file_path, response, category)
                        print(f"Renamed {file_name} to {category}_{response}")
                    else:
                        # Rename the file with just the new name
                        rename_file(file_path, response)
                        print(f"Renamed {file_name} to {response}")
                else:
                    # Do not rename the file if the generated name is not acceptable
                    print(f"Skipping {file_name} due to unacceptable generated name")

if __name__ == "__main__":
    main()
